import pandas as pd
from pptx import Presentation
import win32com.client
import os
import base64
from mailjet_rest import Client
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
SENDER_EMAIL = os.getenv("SENDER_EMAIL")
SENDER_NAME = os.getenv("SENDER_NAME", "Your Name")
MJ_APIKEY_PUBLIC = os.getenv("MJ_APIKEY_PUBLIC")
MJ_APIKEY_PRIVATE = os.getenv("MJ_APIKEY_PRIVATE")
mailjet = Client(auth=(MJ_APIKEY_PUBLIC, MJ_APIKEY_PRIVATE), version='v3.1')

# Read data and set template
df = pd.read_excel("Sample.xlsx")
template_path = "./certificate.pptx"

# Create output directories if they don't exist
os.makedirs("./Files/pptx", exist_ok=True)
os.makedirs("./Files/pdfs", exist_ok=True)

# Initialize PowerPoint once for all conversions
powerpoint = win32com.client.Dispatch("PowerPoint.Application")
powerpoint.Visible = 1

for index, row in df.iterrows():
    
    # Create a copy of the template
    prs = Presentation(template_path)
    
    # Replace placeholders in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if "{{Name}}" in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("{{Name}}", row['Name'])
                if "{{Event}}" in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("{{Event}}", row['Event'])
    
    output_path = f"./Files/pptx/{row['Name']}_{row['Event']}.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    
    # Convert to PDF using PowerPoint COM automation
    pdf_output_path = f"./Files/pdfs/{row['Name']}_{row['Event']}.pdf"
    # Convert to absolute paths
    abs_input_path = os.path.abspath(output_path)
    abs_output_path = os.path.abspath(pdf_output_path)
    
    presentation = powerpoint.Presentations.Open(abs_input_path, WithWindow=False)
    presentation.SaveAs(abs_output_path, 32)  # 32 = ppSaveAsPDF
    presentation.Close()
    print(f"Exported to PDF: {pdf_output_path}")
    
    # Check if email field is empty, skip sending if it is
    if pd.isna(row['Email']) or str(row['Email']).strip() == '':
        print(f"-> Skipped email (no email address provided for {row['Name']})")
        continue
    
    # Send email with Mailjet API
    try:
        # Read and encode PDF file as base64
        with open(pdf_output_path, 'rb') as f:
            pdf_data = f.read()
            pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
        
        # Prepare Mailjet API data
        data = {
            'Messages': [
                {
                    "From": {
                        "Email": SENDER_EMAIL,
                        "Name": SENDER_NAME
                    },
                    "To": [
                        {
                            "Email": row['Email'],
                            "Name": row['Name']
                        }
                    ],
                    "Subject": f"Participation Certificate - {row['Event']} | ANANTYA",
                    "TextPart": f"Dear {row['Name']},\n\nCongratulations on your participation in {row['Event']} at ANANTYA!\n\nWe are delighted to share your participation certificate with you. Please find the certificate attached to this email.\n\nThank you for being a part of ANANTYA and making the event a grand success.\n\nBest regards,\nTeam ANANTYA",
                    "HTMLPart": f"<h2>Dear {row['Name']},</h2><p>Congratulations on your participation in <strong>{row['Event']}</strong> at <strong>ANANTYA</strong>!</p><p>We are delighted to share your participation certificate with you. Please find the certificate attached to this email.</p><p>Thank you for being a part of ANANTYA and making the event a grand success.</p><br/><p>Best regards,<br/><strong>Team ANANTYA</strong></p>",
                    "Attachments": [
                        {
                            "ContentType": "application/pdf",
                            "Filename": f"{row['Name']}_{row['Event']}.pdf",
                            "Base64Content": pdf_base64
                        }
                    ]
                }
            ]
        }
        
        # Send email via Mailjet
        result = mailjet.send.create(data=data)
        
        if result.status_code == 200:
            print(f"-> Successfully sent email with PDF attachment to: {row['Email']}")
        else:
            print(f"-> Failed to send email to {row['Email']}. Status code: {result.status_code}")
            print(f"   Response: {result.json()}")
            
    except Exception as e:
        print(f"-> Failed to send email to {row['Email']}. Error: {e}")

# Close PowerPoint application
powerpoint.Quit()

print("\nAll documents have been processed and emails sent.")

